import fitz
import tempfile
import os
from docx import Document          # pip install python-docx
from pptx import Presentation      # pip install python-pptx


def extract_pdf_text(pdf_path: str) -> list[dict]:
    doc = fitz.open(pdf_path)
    return [
        {"page": page_num + 1, "text": page.get_text()}
        for page_num, page in enumerate(doc)
        if page.get_text().strip()
    ]


def extract_docx_text(path: str) -> list[dict]:
    doc = Document(path)
    return [
        {"page": i + 1, "text": p.text.strip()}
        for i, p in enumerate(doc.paragraphs)
        if p.text.strip()
    ]


def extract_pptx_text(path: str) -> list[dict]:
    prs = Presentation(path)
    pages = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        parts = []

        # Slide title
        if slide.shapes.title and slide.shapes.title.text.strip():
            parts.append(f"[Title] {slide.shapes.title.text.strip()}")

        # All text frames in the slide
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape == slide.shapes.title:
                continue  # already captured above
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    parts.append(text)

        # Notes (speaker notes often contain useful context)
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                parts.append(f"[Notes] {notes_text}")

        if parts:
            pages.append({
                "page": slide_num,
                "text": "\n".join(parts)
            })

    return pages


def extract_file_text(file_bytes: bytes, filename: str) -> list[dict]:
    name = filename.lower()

    if name.endswith(".pdf"):
        suffix = ".pdf"
    elif name.endswith(".docx"):
        suffix = ".docx"
    elif name.endswith(".pptx"):
        suffix = ".pptx"
    else:
        raise ValueError(f"Unsupported file type: {filename}")

    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name

    try:
        if suffix == ".pdf":
            return extract_pdf_text(tmp_path)
        elif suffix == ".docx":
            return extract_docx_text(tmp_path)
        elif suffix == ".pptx":
            return extract_pptx_text(tmp_path)
    finally:
        os.unlink(tmp_path)